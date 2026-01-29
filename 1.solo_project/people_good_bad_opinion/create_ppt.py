"""
임성근 음주운전 논란 여론 분석 - PPT 보고서 생성 스크립트
python-pptx를 이용한 프레젠테이션 자동 생성 (레이아웃 최적화 버전)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
import os

# === 색상 정의 ===
DARK_NAVY = RGBColor(0x1C, 0x28, 0x33)
SLATE_GRAY = RGBColor(0x2E, 0x40, 0x53)
RED = RGBColor(0xE3, 0x37, 0x37)
GREEN = RGBColor(0x00, 0xCC, 0x96)
BLUE = RGBColor(0x63, 0x6E, 0xFA)
ORANGE = RGBColor(0xFF, 0xA1, 0x5A)
PURPLE = RGBColor(0xAB, 0x63, 0xFA)
LIGHT_GREEN = RGBColor(0xB6, 0xE8, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF4, 0xF6, 0xF6)
LIGHT_GRAY = RGBColor(0xAA, 0xB7, 0xB8)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
MEDIUM_GRAY = RGBColor(0x85, 0x92, 0x9E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# === 헬퍼 함수 ===
def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Arial", anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    tf.vertical_anchor = anchor
    return txBox

def add_multiline_text(slide, left, top, width, height, lines, font_size=16,
                       color=DARK_TEXT, bold=False, line_spacing=1.3,
                       alignment=PP_ALIGN.LEFT):
    """lines: list of (text, optional_overrides_dict)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            txt, opts = item
        else:
            txt, opts = item, {}
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(opts.get('size', font_size))
        p.font.color.rgb = opts.get('color', color)
        p.font.bold = opts.get('bold', bold)
        p.font.name = opts.get('font', "Arial")
        p.alignment = opts.get('align', alignment)
        p.space_after = Pt(opts.get('space_after', 6))
        p.line_spacing = line_spacing
    return txBox

def add_metric_card(slide, left, top, width, height, label, value, sub_text=None,
                    value_color=DARK_TEXT, bg_color=WHITE):
    card = add_rounded_rect(slide, left, top, width, height, bg_color)
    # label - Y축 조정
    add_text_box(slide, left + Inches(0.1), top + Inches(0.15), width - Inches(0.2), Inches(0.4),
                 label, font_size=12, color=MEDIUM_GRAY, bold=False, alignment=PP_ALIGN.CENTER)
    # value - 폰트 크기 및 높이 조정
    add_text_box(slide, left + Inches(0.1), top + Incredible(0.55), width - Inches(0.2), Inches(0.6),
                 value, font_size=26, color=value_color, bold=True, alignment=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
    if sub_text:
        add_text_box(slide, left + Inches(0.1), top + height - Inches(0.45), width - Inches(0.2), Inches(0.35),
                     sub_text, font_size=10, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)
    return card

# Incredible is not a function, use Inches
def Incredible(val): return Inches(val)

def add_slide_number(slide, num, total=10):
    add_text_box(slide, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.35),
                 f"{num}/{total}", font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)


# =============================================
# SLIDE 1: 표지
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_shape(slide, 0, 0, prs.slide_width, prs.slide_height, fill_color=DARK_NAVY)
add_shape(slide, 0, Inches(0.8), prs.slide_width, Inches(0.06), fill_color=RED)

add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
             "임성근 음주운전 논란 여론 분석", font_size=46, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.8),
             "데이터 기반 심층 리포트", font_size=24, color=RED, bold=False,
             alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(5), Inches(4.3), Inches(3.333), Inches(0.02), fill_color=LIGHT_GRAY)

add_text_box(slide, Inches(1.5), Inches(4.7), Inches(10), Inches(0.6),
             "백종원 사례 비교 및 사건 소강 주기 분석", font_size=18, color=LIGHT_GRAY,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(6.0), Inches(10), Inches(0.5),
             "2026. 01. 28  |  발표자: Antigravity AI", font_size=14, color=MEDIUM_GRAY,
             alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 1)


# =============================================
# SLIDE 2: 프로젝트 개요
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, prs.slide_width, Inches(1.2), fill_color=DARK_NAVY)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(8), Inches(0.7),
             "프로젝트 개요", font_size=32, color=WHITE, bold=True)

# 분석 배경 박스
add_rounded_rect(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(2.6), OFF_WHITE)
add_text_box(slide, Inches(0.9), Inches(1.75), Inches(5.3), Inches(0.5),
             "분석 배경", font_size=20, color=RED, bold=True)
add_multiline_text(slide, Inches(0.9), Inches(2.3), Inches(5.2), Inches(1.8), [
    ("2026년 1월 19일, 유튜버 임성근 음주운전 자백 사건 발생", {}),
    ("온라인 커뮤니티 및 유튜브 채널 내 부정 여론 급등", {}),
    ("", {'size': 8}),
    ("본 보고서는 해당 사건의 부정 여론 확산 속도를 정밀 분석하고", {}),
    ("백종원 사례와의 비교를 통해 여론이 소강되는 주기를 도출합니다.", {}),
], font_size=13, color=DARK_TEXT)

# 분석 목적 박스
add_rounded_rect(slide, Inches(6.9), Inches(1.6), Inches(5.8), Inches(2.6), OFF_WHITE)
add_text_box(slide, Inches(7.2), Inches(1.75), Inches(5.3), Inches(0.5),
             "분석 목적", font_size=20, color=BLUE, bold=True)
add_multiline_text(slide, Inches(7.2), Inches(2.3), Inches(5.2), Inches(1.8), [
    ("\u2460 부정 여론의 증가율 및 최고점 도달 시점 분석", {'bold': True}),
    ("\u2461 개인 범죄 vs 사업적 논란의 관심도 소요 주기 비교", {'bold': True}),
    ("\u2462 시계열 데이터를 통한 부정 여론의 감쇠 패턴 파악", {'bold': True}),
    ("\u2463 대중의 분노가 소강되는 변곡점 및 주기 예측", {'bold': True}),
], font_size=13, color=DARK_TEXT)

# 분석 데이터 요약
add_rounded_rect(slide, Inches(0.6), Inches(4.6), Inches(12.1), Inches(2.3), DARK_NAVY)
add_text_box(slide, Inches(0.9), Inches(4.75), Inches(3), Inches(0.4),
             "분석 데이터", font_size=18, color=WHITE, bold=True)

cards = [
    ("임성근 댓글", "31,500건+", "논란 발생 직후 집중 수집"),
    ("백종원 댓글", "13,900건+", "논란기 vs 현재 비교군"),
    ("영상 관심지수", "일별 트래킹", "부정 확산 및 소강 지표"),
]
for i, (label, val, sub) in enumerate(cards):
    x = Inches(0.9 + i * 4.05)
    add_metric_card(slide, x, Inches(5.2), Inches(3.6), Inches(1.5),
                    label, val, sub, value_color=RED if i == 0 else BLUE)

add_slide_number(slide, 2)


# =============================================
# SLIDE 3: 분석 방법론 (정제 및 LLM 위주)
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, prs.slide_width, Inches(1.2), fill_color=DARK_NAVY)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(8), Inches(0.7),
             "데이터 분석 방법론", font_size=32, color=WHITE, bold=True)

# 3단계 파이프라인 (BERT 제외)
steps = [
    ("STEP 01", "데이터 수집", "YouTube API v3 활용\n실시간 댓글 데이터 확보", RED),
    ("STEP 02", "텍스트 정제", "Kiwi 형태소 분석 및\n500+ 불용어 자동 필터링", ORANGE),
    ("STEP 03", "정밀 분석 (LLM)", "DeepSeek-V3 엔진 활용\n문맥 기반 6가지 감정 분류", GREEN),
]

for i, (step, title, desc, color) in enumerate(steps):
    x = Inches(1.2 + i * 3.8)
    box = add_rounded_rect(slide, x, Inches(1.6), Inches(3.3), Inches(2.6), OFF_WHITE)
    add_shape(slide, x, Inches(1.6), Inches(3.3), Inches(0.5), fill_color=color)
    add_text_box(slide, x, Inches(1.6), Inches(3.3), Inches(0.5),
                 step, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, x + Inches(0.1), Inches(2.2), Inches(3.1), Inches(0.4),
                 title, font_size=18, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(2.7), Inches(3.1), Inches(1.2),
                 desc, font_size=13, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

    if i < 2:
        arrow_x = x + Inches(3.35)
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, Inches(2.6), Inches(0.35), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = LIGHT_GRAY
        arrow.line.fill.background()

# 감정 카테고리
emotions = [
    ("지지", GREEN), ("분노", RED), ("중립", BLUE), ("실망", ORANGE), ("조롱", PURPLE), ("그외", LIGHT_GREEN)
]
for i, (name, color) in enumerate(emotions):
    x = Inches(1.2 + i * 1.9)
    chip = add_rounded_rect(slide, x, Inches(4.8), Inches(1.7), Inches(1.5), color)
    add_text_box(slide, x, Inches(4.85), Inches(1.7), Inches(0.4),
                 name, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_multiline_text(slide, x + Inches(0.1), Inches(5.3), Inches(1.5), Inches(0.9), [
        ("정밀 분류 모델로", {'size': 9, 'color': WHITE}),
        ("맥락 분석 수행", {'size': 9, 'color': WHITE}),
    ], font_size=10, line_spacing=1.1, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 3)


# =============================================
# SLIDE 4: 핵심 여론 지표
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, prs.slide_width, Inches(1.2), fill_color=DARK_NAVY)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
             "임성근 여론 핵심 지표", font_size=32, color=WHITE, bold=True)

# 지표 카드
metrics = [
    ("총 댓글 수", "31,668건", "유효 데이터", DARK_TEXT),
    ("논란 후 부정 비율", "66.8%", "심각 수준", RED),
    ("심각도 점수", "6.7/10", "보통 수준", RED),
]
for i, (label, val, sub, vc) in enumerate(metrics):
    x = Inches(0.8 + i * 3.1)
    add_metric_card(slide, x, Inches(1.5), Inches(2.8), Inches(1.8), label, val, sub, value_color=vc)

# 도넛 차트
chart_data = CategoryChartData()
chart_data.categories = ['긍정(지지)', '부정(분노/실망/조롱)', '그외']
chart_data.add_series('여론 분포', (50.2, 37.5, 12.2))
cf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(0.5), Inches(3.6), Inches(6), Inches(3.6), chart_data)
plot = cf.chart.plots[0]
plot.has_data_labels = True
plot.data_labels.show_percentage = True
for idx, color in enumerate([RED, GREEN, BLUE]):
    plot.series[0].points[idx].format.fill.solid()
    plot.series[0].points[idx].format.fill.fore_color.rgb = color

# 요약
add_rounded_rect(slide, Inches(7.0), Inches(3.6), Inches(5.8), Inches(3.4), OFF_WHITE)
add_text_box(slide, Inches(7.3), Inches(3.75), Inches(5), Inches(0.4), "핵심 결과 요약", font_size=18, color=RED, bold=True)
add_multiline_text(slide, Inches(7.3), Inches(4.3), Inches(5.2), Inches(2.6), [
    ("\u25CF 부정 여론 비율 18.0% → 66.8%로 급등 (+48.8%p)", {'bold': True}),
    ("\u25CF 논란 후 '조롱' 감정이 42.4%로 지배적 형성", {}),
    ("\u25CF 지지 댓글 평균 좋아요 54.1로 핵심 팬덤은 결집 양상", {}),
    ("\u25CF 장기적 브랜드 타격 및 대중의 조롱 밈(Meme) 소비 우려", {}),
], font_size=13, line_spacing=1.4)

add_slide_number(slide, 4)


# =============================================
# SLIDE 5: 논란 전후 감정 변화 (바 차트)
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, prs.slide_width, Inches(1.2), fill_color=DARK_NAVY)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
             "사건 전후 감정 격차 분석", font_size=32, color=WHITE, bold=True)

chart_data_cmp = CategoryChartData()
chart_data_cmp.categories = ['지지', '분노', '중립', '실망', '조롱', '그외']
chart_data_cmp.add_series('논란 전', (50.6, 1.4, 22.4, 1.7, 14.9, 9.0))
chart_data_cmp.add_series('논란 후', (9.46, 15.0, 6.65, 9.31, 42.4, 17.1))
cf_cmp = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.8), Inches(1.6), Inches(8), Inches(5), chart_data_cmp)
cf_cmp.chart.series[0].format.fill.solid()
cf_cmp.chart.series[0].format.fill.fore_color.rgb = BLUE
cf_cmp.chart.series[1].format.fill.solid()
cf_cmp.chart.series[1].format.fill.fore_color.rgb = RED

add_rounded_rect(slide, Inches(9.2), Inches(1.6), Inches(3.6), Inches(5.0), OFF_WHITE)
add_text_box(slide, Inches(9.4), Inches(1.8), Inches(3.2), Inches(0.5), "관전 포인트", font_size=18, color=RED, bold=True)
add_multiline_text(slide, Inches(9.4), Inches(2.5), Inches(3.2), Inches(4.0), [
    ("1. 조롱 감정의 폭발적 증가 (14.9% → 42.4%)", {'bold': True}),
    ("2. 충성 팬층을 의미하는 '지지'의 80% 이상 소멸", {}),
    ("3. 사건의 도덕성 문제로 '분노' 비율 약 10배 증가", {}),
    ("4. '중립' 성향이 줄고 냉소적인 여론으로 결집", {}),
], font_size=13, line_spacing=1.6)

add_slide_number(slide, 5)


# =============================================
# SLIDE 6~8 (간단히 유지)
# =============================================
# Slide 6: 백종원 비교
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, prs.slide_width, Inches(1.2), fill_color=DARK_NAVY)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7), "백종원 vs 임성근 비교 분석", font_size=32, color=WHITE, bold=True)

# 차트 데이터 (가로 막대형)
c_data_baek = CategoryChartData()
c_data_baek.categories = ['조롱', '분노', '그외', '실망', '지지', '중립']
c_data_baek.add_series('임성근(논란후)', (42.4, 15.0, 17.1, 9.31, 9.46, 6.65))
c_data_baek.add_series('백종원', (38.0, 24.7, 14.0, 10.5, 7.14, 5.59))
cf_baek = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.8), Inches(1.6), Inches(6.5), Inches(5), c_data_baek)
cf_baek.chart.series[0].format.fill.solid()
cf_baek.chart.series[0].format.fill.fore_color.rgb = RED
cf_baek.chart.series[1].format.fill.solid()
cf_baek.chart.series[1].format.fill.fore_color.rgb = BLUE

add_rounded_rect(slide, Inches(7.5), Inches(1.6), Inches(5.3), Inches(5.0), OFF_WHITE)
add_text_box(slide, Inches(7.7), Inches(1.8), Inches(5), Inches(0.5), "핵심 차이점 및 가설 검증", font_size=18, color=DARK_NAVY, bold=True)
add_multiline_text(slide, Inches(7.7), Inches(2.5), Inches(4.8), Inches(4.0), [
    ("■ 백종원: 조롱 38.0%, 분노 24.7%", {'bold': True}),
    ("- 사업적 실책성 논란 위주 -> 공격적 성향(분노)이 상대적으로 강함", {'size': 12}),
    ("■ 임성근: 조롱 42.4%, 분노 15.0%", {'bold': True}),
    ("- 개인적/도덕적 실망 -> 냉소적 반응(조롱)이 지배적으로 형성", {'size': 12}),
    ("", {'size': 8}),
    ("■ 가설 검증 결과", {'bold': True, 'color': RED}),
    ("- 두 사례 모두 '조롱'이 '분노'를 압도하여 여론을 주도함", {'size': 12}),
    ("- 임성근 사례에서 감정 격차(27.4%p)가 백종원(13.3%p)보다 약 2배 큼", {'size': 12}),
], font_size=13, line_spacing=1.4)
add_slide_number(slide, 6)

# Slide 7: 시계열
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, prs.slide_width, Inches(1.2), fill_color=DARK_NAVY)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7), "사례별 시계열 부정 여론 추이", font_size=32, color=WHITE, bold=True)

# 임성근 주간 추이 (좌측)
c_data_lim = CategoryChartData()
c_data_lim.categories = ['1/13', '1/14', '1/15', '1/16', '1/17', '1/18', '1/19', '1/20', '1/21', '1/22']
c_data_lim.add_series('임성근 부정 비율(%)', (0.0, 9.3, 4.8, 11.6, 7.9, 22.7, 72.2, 72.2, 67.0, 70.1))
cf_lim = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.5), c_data_lim)
cf_lim.chart.series[0].format.line.color.rgb = RED
cf_lim.chart.series[0].format.line.width = Pt(3)

# 백종원 월별 추이 (우측)
c_data_baek_ts = CategoryChartData()
c_data_baek_ts.categories = ['25.04', '25.05', '25.06', '25.07', '25.08', '25.09', '25.10']
c_data_baek_ts.add_series('백종원 부정 비율(%)', (85.1, 80.2, 57.0, 76.7, 51.7, 60.7, 46.0))
cf_baek_ts = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(7.0), Inches(1.8), Inches(5.8), Inches(4.5), c_data_baek_ts)
cf_baek_ts.chart.series[0].format.line.color.rgb = BLUE
cf_baek_ts.chart.series[0].format.line.width = Pt(3)
add_slide_number(slide, 7)

# Slide 8: 관심지수
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, prs.slide_width, Inches(1.2), fill_color=DARK_NAVY)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7), "영상 관심지수 트렌드 분석", font_size=32, color=WHITE, bold=True)
add_metric_card(slide, Inches(0.8), Inches(1.5), Inches(2.8), Inches(1.5), "총 영상 수", "1,655건", "↓ 22개 (수렴 중)", SLATE_GRAY)
add_metric_card(slide, Inches(3.9), Inches(1.5), Inches(2.8), Inches(1.5), "총 조회수", "4.13억", "↓ 51만 (관심도 소강)", SLATE_GRAY)
add_metric_card(slide, Inches(7.0), Inches(1.5), Inches(2.8), Inches(1.5), "총 좋아요", "391만", "↓ 1,314 (추진력 감소)", SLATE_GRAY)
add_metric_card(slide, Inches(10.1), Inches(1.5), Inches(2.8), Inches(1.5), "총 댓글 수", "21.4만", "↑ 15건 (여론 고착화)", SLATE_GRAY)
add_slide_number(slide, 8)


# =============================================
# SLIDE 9: 결론 및 인사이트 (중요 수정)
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, prs.slide_width, Inches(1.2), fill_color=DARK_NAVY)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
             "핵심 인사이트 & 향후 전망", font_size=32, color=WHITE, bold=True)

insight_cards = [
    {
        "title": "논란의 심각성과 성격", "color": RED,
        "bullets": [
            "부정 비율 18% → 66.8%로 급격 팽창",
            "단순 분노를 넘어 '조롱'(42.4%)이 핵심 기조",
            "초기 분노가 가라앉아도 냉소적 여론이 고착화될 위험",
            "대중의 밈(Meme) 소비로 인한 지속적 이미지 타격 발생"
        ]
    },
    {
        "title": "부정 여론 증가 및 소강 주기", "color": BLUE,
        "bullets": [
            "백종원 사례: 약 3~4개월 후 관심도 50% 수준 소강",
            "임성근: 도덕적 결함으로 인해 더 긴 소강 주기 예상",
            "1,655개 영상의 도달 범위가 넓어 소강까지 12개월 이상 소요",
            "과열 시기 이후 발생하는 '무관심' 변곡점이 진정한 소강 상태"
        ]
    },
    {
        "title": "데이터 기반 대응 전략", "color": GREEN,
        "bullets": [
            "감정 변곡점 포착을 위한 '댓글 톤(Tone)' 상시 트래킹",
            "지지층(평균 좋아요 54.1) 결집 및 중립층 수용도 개선 필요",
            "부정 확산 속도 제어를 없는 진정성 있는 소통 전략 필수",
            "데이터는 단순 비난을 넘어 사회적 공감의 한계를 증명"
        ]
    }
]

for i, item in enumerate(insight_cards):
    x = Inches(0.8 + i * 4.1)
    # 카드 배경
    add_rounded_rect(slide, x, Inches(1.6), Inches(3.7), Inches(5.2), OFF_WHITE)
    # 헤더
    header = add_shape(slide, x, Inches(1.6), Inches(3.7), Inches(0.6), fill_color=item['color'])
    add_text_box(slide, x, Inches(1.6), Inches(3.7), Inches(0.6), item['title'],
                 font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    
    # 불릿 포인트 - add_multiline_text 활용하여 하나의 박스에 담음 (겹침 방지)
    bullet_lines = [(f"\u2022 {b}", {}) for b in item['bullets']]
    add_multiline_text(slide, x + Inches(0.2), Inches(2.4), Inches(3.3), Inches(4.0),
                       bullet_lines, font_size=12, line_spacing=1.3)

add_slide_number(slide, 9)


# =============================================
# SLIDE 10: Q&A
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, prs.slide_width, prs.slide_height, fill_color=DARK_NAVY)
add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10.333), Inches(1.2),
             "Q & A", font_size=60, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10.333), Inches(0.6),
             "경청해 주셔서 감사합니다.", font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 10)


# === 저장 ===
output_path = os.path.join(os.path.dirname(__file__), "임성근_음주운전_논란_여론분석_보고서.pptx")
prs.save(output_path)
print(f"PPT 생성 완료: {output_path}")
